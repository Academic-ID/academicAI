import sys
import os
sys.path.insert(0, os.path.abspath(os.path.join(os.getcwd(), os.pardir)))
from pptx import Presentation
from docx import Document
import openpyxl
from pdfminer.high_level import extract_text as pdfminer_extract
from pdf2image import convert_from_path
import tiktoken
from striprtf.striprtf import rtf_to_text
from client import client, parse_JSON, default_model, vision_model
from io import BytesIO
import base64
import time
from functools import wraps

def retry(max_retries=3, delay=2):
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            retries = 0
            while retries < max_retries:
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    retries += 1
                    if retries < max_retries:
                        print(f"Retry {retries}/{max_retries} for {func.__name__} due to: {str(e)}")
                        time.sleep(delay * retries)
                    else:
                        raise e
        return wrapper
    return decorator

def pdf_to_base64_images(pdf_path):
    if not pdf_path.lower().endswith('.pdf'):
        return 'Not a PDF file'
    """
    Converts each page of a PDF to a base64 encoded image.

    Args:
    pdf_path (str): The file path to the PDF.

    Returns:
    list: A list of base64 encoded strings for each page of the PDF.
    """
    # Convert PDF to a list of image objects
    images = convert_from_path(pdf_path)    

    base64_images = []    
    for image in images:
        # Convert image to bytes using an in-memory buffer
        buffer = BytesIO()
        image.save(buffer, format="JPEG")  # You might need to adjust the format as necessary
        # Encode the bytes to base64 and append to the list
        base64_images.append(base64.b64encode(buffer.getvalue()).decode('utf-8'))

    return base64_images

@retry()
async def get_img_chat_completion(messages, max_res_tokens=500, temp=0.1):
    try:
        response = client.chat.completions.create(
            model=vision_model,
            messages=messages,
            max_tokens=max_res_tokens,
            temperature=temp
        )
        return response.choices[0].message.content                
    except Exception as e:        
        raise Exception(f'Error witth vision request: {e}')

async def get_summary_vision(doc_path, user_instruction = ''):    
    pdf_imgs = pdf_to_base64_images(doc_path)    
    if pdf_imgs == 'Not a PDF file':
        text_summary = "There was an error retrieving a summary for this document"
        return text_summary
    system_prompt="Large language model is a highly-paid expert in PDF document summarisation. Large language model takes time to read documents and consider the meaning of the entire document before providing a fully thought out response. The user will send you some images of a document of thiers to be summarised. Please provide a summary of the document. The summary should be detailed, but needs to be concise and use plain english. The summary should not include the document name. The summary is being used for a semantic database so maintaining the semantic meaning of the text is paramount."
    content_objs = []
    if user_instruction != '':
        content_objs.append({"type": "text", "text": user_instruction})
    for img in pdf_imgs[:9]:
        content_objs.append({"type": "image_url", "image_url": f"data:image/jpeg;base64,{img}"})
    messages = [{"role": "system", "content": system_prompt}, {"role": "user", "content": content_objs}]
    response = await get_img_chat_completion(messages)    
    return response

@retry()
async def get_summary(doc_text, doc_type, user_instruction = ''):
    """Returns a json object with is_text_present and summary keys. is_text_present is a boolean value indicating whether the text is present and readable. summary is the summary of the document."""
    system_prompt=f"Large language model is a highly-paid expert in document summarisation - specifically {doc_type} files . Large language model takes time to read documents and consider the meaning of the entire document before providing a fully thought out response."
    sum_fn = [
    {
      'type': 'function',
      'function': {
        'name': 'add_doc_summary_to_semantic_database',
        'description': "This function requires the assistant to provide two parameters. First, the assistant must determine if the provided text is not empty and is legible enough to provide a summary. If it is, the second parameter should include that summary. If the text is not present or not legible, the summary parameter should be returned as 'N/A'",
        'parameters': {
          'type': 'object',
          'properties': {
            'is_text_present': {
              'type': 'boolean',
              'description':
                'The text is being extracted from documents include PDF files. There are occasions where the text extraction fails and no content is returned or the content is nonsensical. If the content is present, return True, if not present (or text is not capable of being summarised), return False.',
            },
            'summary': {
              'type': 'string',
              'description':
                'The summary of the document. If the text is not present or not legible, return "N/A". The summary should be detailed, but needs to be concise and use plain english. The summary should not include the document name. The summary is being used for a semantic database so maintaining the semantic meaning of the text is paramount.',
            },
          },
          'required': ['is_text_present', 'summary'],
        },
      },
    }]
    response = client.chat.completions.create(
                    model=default_model,
                    messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": f'{user_instruction}{doc_text}'}],
                    max_tokens=4096,
                    temperature=0.1,
                    tools=sum_fn,
                    tool_choice={ 'type': 'function', 'function': { 'name': 'add_doc_summary_to_semantic_database' } },        
                )
    response_message = response.choices[0].message
    tool_call_args = response_message.tool_calls[0].function.arguments
    return parse_JSON(tool_call_args)

def extract_text(file_path):
    if file_path.endswith('.pdf'):
        text = pdfminer_extract(file_path)
        # Splitting text into lines
        lines = text.split('\n')
        # Joining lines if a line starts with a lowercase letter
        joined_lines = []
        for line in lines:
            if line and line[0].islower():
                if joined_lines:  # ensure joined_lines is not empty
                    joined_lines[-1] = joined_lines[-1] + line
                else:  # if joined_lines is empty, just append the line
                    joined_lines.append(line)
            else:
                joined_lines.append(line)
        # Removing newline characters which might be in the middle of the words
        joined_lines = [line.replace('\n', '') for line in joined_lines]
        # Removing consecutive duplicates
        non_repeat_lines = [joined_lines[i] for i in range(
            len(joined_lines)) if i == 0 or joined_lines[i] != joined_lines[i-1]]
        return '\n'.join(non_repeat_lines)
    elif file_path.endswith('.docx'):
        doc = Document(file_path)
        paragraphs = doc.paragraphs
        non_repeat_paras = [paragraphs[i].text for i in range(
            len(paragraphs)) if i == 0 or paragraphs[i].text != paragraphs[i-1].text]
        return '\n'.join(non_repeat_paras)
    elif file_path.endswith('.pptx'):
        prs = Presentation(file_path)
        text_shapes = [
            shape for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        non_repeat_shapes = [text_shapes[i].text for i in range(
            len(text_shapes)) if i == 0 or text_shapes[i].text != text_shapes[i-1].text]
        return '\n'.join(non_repeat_shapes)
    elif file_path.endswith('.xlsx'):
        workbook = openpyxl.load_workbook(file_path)
        sheet = workbook.active
        text = ""
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                text += str(cell) + " "
            text += "\n"
        return text
    elif file_path.endswith('.txt'):
        with open(file_path, 'r') as f:
            text = f.read()
        return text
    elif file_path.endswith('.rtf'):
        with open(file_path, 'r') as infile:
            content = infile.read()
            text = rtf_to_text(content)
        return text
    else:
        raise ValueError('Unsupported file type')

async def process_doc(doc, user_instruction, pbar):    
    # plain function to extract text from doc
    doc_text = extract_text(doc["path"])
    doc_type = doc["path"].split('.')[-1]

    # Tokenize the input string.
    encoding = tiktoken.get_encoding("cl100k_base")
    num_tokens = len(encoding.encode(doc_text))    

    # reduce string size for summary generation
    # 100000 token should be enough to get a general overview of the content of a document
    first_50000_tokens = doc_text
    if num_tokens > 100000:        
        first_50000_tokens = encoding.encode(doc_text)[:100000]
        first_50000_tokens = encoding.decode(first_50000_tokens)

    # get summary of document
    text_summary = ''
    try:
        failed_responses = {'n/a', 'false'}
        summary_obj = await get_summary(first_50000_tokens, doc_type, user_instruction)        
        summary_str = str(summary_obj["summary"]).lower().strip("'")  
        if (len(summary_str) < 10 or summary_str in failed_responses or not summary_obj["is_text_present"] or str(summary_obj["is_text_present"]).lower().strip("'") in failed_responses):
            try:
                text_summary = await get_summary_vision(doc["path"])
            except Exception as e:
                text_summary = f"There was an error retrieving a summary for this document. {e}"
        else:
            text_summary = summary_obj["summary"]
        
    except Exception as e:
        text_summary = f"There was an error retrieving a summary for this document. {e}"    
    pbar.update(1)
    return text_summary
